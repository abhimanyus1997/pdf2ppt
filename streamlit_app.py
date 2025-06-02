import os
import fitz  # PyMuPDF
import pytesseract
import tempfile
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
import streamlit as st

DPI = 150
MAX_SLIDE_INCH = 56

def pdf_to_pptx_and_images(pdf_file, enable_ocr=False):
    with tempfile.TemporaryDirectory() as temp_dir:
        pdf_path = os.path.join(temp_dir, "input.pdf")
        with open(pdf_path, "wb") as f:
            f.write(pdf_file.read())

        doc = fitz.open(pdf_path)
        total_pages = len(doc)
        prs = Presentation()

        slide_images = []

        for i in range(total_pages):
            page = doc.load_page(i)
            pix = page.get_pixmap(dpi=DPI)
            img_path = os.path.join(temp_dir, f"page_{i}.png")
            pix.save(img_path)

            width_in = min(pix.width / DPI, MAX_SLIDE_INCH)
            height_in = min(pix.height / DPI, MAX_SLIDE_INCH)
            prs.slide_width = Inches(width_in)
            prs.slide_height = Inches(height_in)

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

            if enable_ocr:
                image = Image.open(img_path)
                ocr_text = pytesseract.image_to_string(image)
                if ocr_text.strip():
                    left = Inches(0.2)
                    top = Inches(0.2)
                    width = prs.slide_width - Inches(0.4)
                    height = Inches(1.5)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    tf = text_box.text_frame
                    tf.text = ocr_text.strip()[:1000]
                    tf.paragraphs[0].font.size = Pt(12)

            # Collect slide images for preview
            slide_images.append(img_path)

        pptx_path = os.path.join(temp_dir, "output.pptx")
        prs.save(pptx_path)
        with open(pptx_path, "rb") as f:
            pptx_bytes = f.read()

        return pptx_bytes, slide_images

# Streamlit UI
st.set_page_config(page_title="📄 PDF to PPTX Converter", layout="centered")
st.title("📄 PDF to PPTX Converter")
st.markdown("Convert your PDF into a PowerPoint presentation. Optionally overlay OCR text.")

uploaded_pdf = st.file_uploader("Upload a PDF file", type=["pdf"])
enable_ocr = st.checkbox("Enable OCR text overlay", value=False)

if uploaded_pdf and st.button("Convert to PPTX"):
    with st.spinner("⏳ Converting..."):
        pptx_bytes, slide_images = pdf_to_pptx_and_images(uploaded_pdf, enable_ocr)

        st.success("✅ Conversion complete!")

        # Preview PDF pages / PPT slides as images
        st.subheader("Preview slides")
        for img_path in slide_images:
            img = Image.open(img_path)
            st.image(img, use_column_width=True)

        st.download_button("📥 Download PPTX", pptx_bytes, file_name="converted.pptx")
